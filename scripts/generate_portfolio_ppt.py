"""포트폴리오용 PPT 자동 생성 (python-pptx + matplotlib).

전력수급 실시간 모니터링·다층 이상탐지 산출물을 면접/포트폴리오용 슬라이드 덱으로 묶는다.
차트(부하곡선·발전믹스·이상탐지 타임라인·3계층 비교)는 matplotlib로 그려 PNG로 임베드한다.

데이터가 없으면 합성 시계열로 데모 덱을 생성한다 (API 키 등록 전에도 산출물 확인 가능).

실행:
    .venv/Scripts/python.exe scripts/generate_portfolio_ppt.py
출력:
    reports/portfolio/YYYY-MM-DD.pptx
    reports/portfolio/latest.pptx
"""
from __future__ import annotations

import os
import sys
import tempfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Optional

_PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

import matplotlib  # noqa: E402

matplotlib.use("Agg")  # headless (CI)
import matplotlib.font_manager as fm  # noqa: E402
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
import pandas as pd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from flows.analysis_flow import detection_summary, run_analysis  # noqa: E402
from src.storage import database  # noqa: E402

KST = timezone(timedelta(hours=9))
REPORTS_DIR = _PROJECT_ROOT / "reports" / "portfolio"

# 색상 팔레트
_NAVY = RGBColor(0x14, 0x28, 0x50)
_BLUE = RGBColor(0x1F, 0x77, 0xB4)
_GRAY = RGBColor(0x55, 0x55, 0x55)
_RED = RGBColor(0xD6, 0x27, 0x28)

_FONT_CANDIDATES = [
    os.environ.get("PORTFOLIO_FONT", ""),
    "C:/Windows/Fonts/malgun.ttf",
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
]

SOURCE_COLORS = {
    "원자력": "#4e79a7", "LNG": "#f28e2b", "석탄": "#59a14f",
    "신재생": "#76b7b2", "수력": "#edc948", "양수": "#b07aa1", "유류": "#ff9da7",
}


def _setup_korean_font() -> Optional[str]:
    for c in _FONT_CANDIDATES:
        if c and Path(c).exists():
            try:
                fm.fontManager.addfont(c)
                name = fm.FontProperties(fname=c).get_name()
                plt.rcParams["font.family"] = name
                plt.rcParams["axes.unicode_minus"] = False
                return name
            except Exception:
                continue
    print("⚠️ 한글 폰트를 찾지 못했습니다. 차트 한글이 깨질 수 있습니다.")
    return None


# ---------------------------------------------------------------------------
# 합성 데이터 (데이터 부족 시 데모용)
# ---------------------------------------------------------------------------
def _make_synthetic() -> pd.DataFrame:
    """일주기·주간주기 + 노이즈 + 이상 이벤트가 섞인 합성 전력수급 시계열."""
    n = 288 * 7  # 1주
    idx = pd.date_range(
        datetime.now(KST).replace(tzinfo=None) - timedelta(days=7),
        periods=n, freq="5min",
    )
    t = np.arange(n)
    rng = np.random.default_rng(42)
    # 일주기(낮 피크) + 주간(주말 저하) + 노이즈
    daily = 10000 * np.sin(2 * np.pi * (t % 288) / 288 - np.pi / 2)
    weekly = -3000 * ((idx.dayofweek >= 5).astype(float))
    load = 65000 + daily + weekly + rng.normal(0, 800, n)
    # 이상 이벤트: 급증 스파이크 2회
    load[int(n * 0.4)] += 18000
    load[int(n * 0.75)] += 15000
    supply = load + rng.uniform(8000, 15000, n)
    reserve_rate = (supply - load) / load * 100
    return pd.DataFrame({
        "ts": idx,
        "supply_capacity": supply,
        "current_load": load,
        "forecast_load": load + rng.normal(0, 1500, n),
        "reserve_power": supply - load,
        "reserve_rate": reserve_rate,
        "oper_reserve_power": (supply - load) * 0.7,
        "oper_reserve_rate": reserve_rate * 0.7,
    })


# ---------------------------------------------------------------------------
# 차트 (matplotlib → PNG)
# ---------------------------------------------------------------------------
def _chart_load_curve(clean: pd.DataFrame, tmpdir: Path) -> Path:
    """부하 곡선 — 현재수요·공급능력·예측수요."""
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=150)
    df = clean.set_index("ts").tail(288 * 3)  # 최근 3일
    ax.plot(df.index, df["current_load"], label="현재수요", color="#1f77b4", lw=1.5)
    ax.plot(df.index, df["supply_capacity"], label="공급능력", color="#2ca02c", lw=1, ls="--")
    ax.plot(df.index, df["forecast_load"], label="예측수요", color="#aec7e8", lw=1, ls=":")
    ax.set_title("전력 부하 곡선 (제조 생산라인 throughput 유사)", fontsize=13, fontweight="bold")
    ax.set_ylabel("MW"); ax.legend(fontsize=9); ax.grid(alpha=0.3)
    fig.autofmt_xdate(); fig.tight_layout()
    path = tmpdir / "load.png"; fig.savefig(path); plt.close(fig)
    return path


def _chart_reserve(clean: pd.DataFrame, tmpdir: Path) -> Path:
    """예비율 추이 + 경보 임계선."""
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=150)
    df = clean.set_index("ts").tail(288 * 3)
    ax.plot(df.index, df["reserve_rate"], label="공급예비율", color="#ff7f0e", lw=1.5)
    ax.fill_between(df.index, df["reserve_rate"], alpha=0.1, color="#ff7f0e")
    ax.axhline(10, color="orange", ls="--", lw=1, label="주의 10%")
    ax.axhline(5, color="red", ls="--", lw=1, label="심각 5%")
    ax.set_title("공급예비율 추이 (안전재고·여유율 유사)", fontsize=13, fontweight="bold")
    ax.set_ylabel("%"); ax.legend(fontsize=9); ax.grid(alpha=0.3)
    fig.autofmt_xdate(); fig.tight_layout()
    path = tmpdir / "reserve.png"; fig.savefig(path); plt.close(fig)
    return path


def _chart_anomaly_timeline(result: dict, tmpdir: Path) -> Path:
    """이상탐지 타임라인 — 원시 시계열 + L1/L2 이상 마킹."""
    clean = result["clean_df"]
    target = result["target_col"]
    df = clean.set_index("ts")
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=150)
    ax.plot(df.index, df[target], color="#1f77b4", lw=1.2, label=target)

    det = result.get("detections", {})
    if "L1_ewma" in det and "df" in det["L1_ewma"]:
        e = det["L1_ewma"]["df"]
        anom = e[e["anomaly"]]
        if not anom.empty:
            ax.scatter(anom.index, anom["value"], color="red", s=40,
                       marker="x", label=f"L1 EWMA ({len(anom)})", zorder=5)
    if "L2_iforest" in det and "df" in det["L2_iforest"]:
        l2 = det["L2_iforest"]["df"].set_index("ts")
        anom = l2[l2["anomaly"] == True]
        if not anom.empty:
            ax.scatter(anom.index, anom[target], facecolors="none",
                       edgecolors="orange", s=80, label=f"L2 IForest ({len(anom)})", zorder=4)
    ax.set_title("다층 이상탐지 타임라인 (공정 이상 감지)", fontsize=13, fontweight="bold")
    ax.set_ylabel(target); ax.legend(fontsize=9); ax.grid(alpha=0.3)
    fig.autofmt_xdate(); fig.tight_layout()
    path = tmpdir / "anomaly.png"; fig.savefig(path); plt.close(fig)
    return path


def _chart_layer_compare(result: dict, tmpdir: Path) -> Path:
    """3계층 탐지 건수 비교 막대."""
    det = result.get("detections", {})
    labels, counts = [], []
    for key, label in [("L1_ewma", "L1 EWMA"), ("L1_cusum", "L1 CUSUM"),
                       ("L2_iforest", "L2 IForest"), ("residual", "잔차기반"),
                       ("L3_lstm", "L3 LSTM-AE")]:
        d = det.get(key, {})
        if not d.get("skipped") and "n_anomalies" in d:
            labels.append(label); counts.append(d["n_anomalies"])
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=150)
    colors = ["#1f77b4", "#aec7e8", "#ff7f0e", "#2ca02c", "#d62728"][:len(labels)]
    ax.bar(labels, counts, color=colors)
    for i, c in enumerate(counts):
        ax.text(i, c, str(c), ha="center", va="bottom", fontsize=10)
    ax.set_title("계층별 이상 감지 건수 비교", fontsize=13, fontweight="bold")
    ax.set_ylabel("감지 건수"); ax.grid(alpha=0.3, axis="y")
    fig.tight_layout()
    path = tmpdir / "compare.png"; fig.savefig(path); plt.close(fig)
    return path


def _chart_genmix(tmpdir: Path) -> Optional[Path]:
    """발전믹스 파이 (데이터 있을 때만)."""
    gen = database.load_generation_df()
    if gen.empty:
        return None
    latest_ts = gen["ts"].max()
    snap = gen[gen["ts"] == latest_ts].groupby("source")["generation_mw"].sum()
    fig, ax = plt.subplots(figsize=(7, 4.5), dpi=150)
    colors = [SOURCE_COLORS.get(s, "#ccc") for s in snap.index]
    ax.pie(snap.values, labels=snap.index, colors=colors, autopct="%1.1f%%",
           startangle=90, wedgeprops={"width": 0.5})
    ax.set_title(f"발전믹스 ({latest_ts:%m-%d %H:%M})", fontsize=13, fontweight="bold")
    fig.tight_layout()
    path = tmpdir / "genmix.png"; fig.savefig(path); plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# 슬라이드 헬퍼
# ---------------------------------------------------------------------------
def _txt(slide, left, top, width, height, text, size=14, bold=False,
         color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return box


def _heading(slide, text):
    _txt(slide, 0.6, 0.4, 12.1, 0.9, text, size=24, bold=True, color=_NAVY)


def _bullets(slide, items, size=15):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(10)
        run = para.add_run(); run.text = f"•  {item}"
        run.font.size = Pt(size); run.font.color.rgb = _GRAY


def _add_chart_slide(prs, title, img, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, title)
    slide.shapes.add_picture(str(img), Inches(1.2), Inches(1.5), width=Inches(10.8))
    _txt(slide, 0.6, 6.7, 12.1, 0.7, caption, size=12, color=_GRAY)


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------
def build_ppt(result: dict, today: datetime, is_synthetic: bool) -> Path:
    _setup_korean_font()
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)

    clean = result["clean_df"]
    eda = result["eda"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 타이틀
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _txt(slide, 0.6, 1.8, 12.1, 1.3, "실시간 전력수급 이상탐지 모니터링",
         size=34, bold=True, color=_NAVY)
    _txt(slide, 0.6, 3.0, 12.1, 0.8,
         "통계(EWMA/CUSUM) · ML(Isolation Forest) · 딥러닝(LSTM-AutoEncoder) 다층 이상탐지",
         size=17, color=_BLUE)
    tr = eda.get("time_range")
    rng = f"{tr[0]:%Y-%m-%d %H:%M} ~ {tr[1]:%Y-%m-%d %H:%M}" if tr else "—"
    demo = "  ※ (현재 합성 데모 데이터 — KPX API 키 등록 후 실데이터로 자동 교체)" if is_synthetic else ""
    _txt(slide, 0.6, 4.3, 12.1, 1.6,
         f"분석 기간: {rng}\n"
         f"데이터: 총 {eda['n_rows']:,}건 · 5분 해상도 · 전국 계통 기준\n"
         f"생성일: {today:%Y-%m-%d %H:%M KST}{demo}",
         size=14, color=_GRAY)

    # 2. 개요 — 제조 메타포
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "프로젝트 개요 — 전력계통을 제조 생산라인에 매핑")
    _bullets(slide, [
        "현재 전력수요(부하) = 생산라인 throughput / 공급예비율 = 안전재고·여유율",
        "발전원별 발전량(원자력·LNG·석탄·신재생) = 다중 설비 가동 상태",
        "예비율 임계선 = 규격 한계(USL/LSL) / 수요 급변·예비율 급락 = 공정 이상",
        "KPX 5분 단위 OpenAPI 자동 수집 → SQLite → GitHub Actions 무중단 운영",
        "차별점: SPC 중심의 자매 프로젝트와 달리 ML/DL 기반 이상탐지로 기법 차별화",
    ])

    # 3. 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "워크플로우 — 수집→전처리→다층탐지→대시보드/PPT")
    _bullets(slide, [
        "수집: collect_flow.py가 KPX 수급/발전믹스 5분 폴링 → 멱등 upsert",
        "전처리(EDA): preprocess.py — 리샘플·결측보간·파생변수(변화율·예측오차)",
        "이상탐지: analysis_flow.py가 L1 통계→L2 ML→L3 DL 순차 게이트 실행",
        "대시보드: Streamlit 6페이지 (모니터링·이상타임라인·발전믹스·탐지비교·수요예측·지도)",
        "포트폴리오: 매 분석 후 PPT 자동 생성 + 데일리 리포트",
        "운영: GitHub Actions cron 무중단 수집 — 컴퓨터를 꺼도 데이터 누적",
    ])

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        _add_chart_slide(prs, "모니터링 — 부하 곡선", _chart_load_curve(clean, tmp),
                         "현재수요·공급능력·예측수요를 5분 해상도로 추적. 일주기 패턴이 뚜렷.")
        _add_chart_slide(prs, "모니터링 — 공급예비율", _chart_reserve(clean, tmp),
                         "예비율이 임계선(주의 10%·심각 5%)에 접근하면 블랙아웃 위험 — 경보 발령.")
        _add_chart_slide(prs, "이상탐지 — 다층 타임라인", _chart_anomaly_timeline(result, tmp),
                         "동일 시계열에 L1(통계)·L2(ML) 탐지 결과를 중첩. 급변 구간을 자동 포착.")
        _add_chart_slide(prs, "이상탐지 — 계층별 비교", _chart_layer_compare(result, tmp),
                         "계층별 감지 건수 비교. 단순 임계값 대비 다층 탐지로 거짓경보·누락 균형.")
        genmix = _chart_genmix(tmp)
        if genmix:
            _add_chart_slide(prs, "발전믹스 — 발전원별 구성", genmix,
                             "원자력·LNG·석탄·신재생 발전 비중. 설비 가동 포트폴리오 모니터링.")

        # 결론
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _heading(slide, "핵심 성과 & 인사이트")
        bullets = [
            f"총 {eda['n_rows']:,}건 · {eda.get('duration_hours', 0):.0f}시간 분량 5분 해상도 수집·분석",
            f"다층 이상탐지 결과: {detection_summary(result)}",
        ]
        if "load_peak" in eda:
            bullets.append(
                f"최대 부하: {eda['load_peak']['value']:,.0f} MW "
                f"({eda['load_peak']['ts']:%m-%d %H:%M})"
            )
        if "min_reserve" in eda:
            bullets.append(
                f"최저 예비율: {eda['min_reserve']['value']:.1f}% "
                f"({eda['min_reserve']['ts']:%m-%d %H:%M}) — 위험 시점"
            )
        bullets.append("전처리→탐지→리포트 전 과정 자동화 — 무중단·무비용 파이프라인")
        _bullets(slide, bullets)

        dated = REPORTS_DIR / f"{today:%Y-%m-%d}.pptx"
        latest = REPORTS_DIR / "latest.pptx"
        prs.save(str(dated)); prs.save(str(latest))

    return dated


def main() -> None:
    today = datetime.now(KST)
    df = database.load_df()
    is_synthetic = df.empty
    if is_synthetic:
        print("⚠️ 수집 데이터 없음 — 합성 데모 데이터로 PPT 생성")
        df = _make_synthetic()

    result = run_analysis(df)
    if result.get("empty"):
        print("분석 실패 — 데이터 없음")
        return

    path = build_ppt(result, today, is_synthetic)
    n_slides = len(Presentation(str(path)).slides._sldIdLst)
    print(f"✅ 포트폴리오 PPT 생성: {path.relative_to(_PROJECT_ROOT)}")
    print(f"   슬라이드 {n_slides}장 · 데이터 {result['n_rows']:,}건")
    print(f"   탐지 요약: {detection_summary(result)}")


if __name__ == "__main__":
    main()
